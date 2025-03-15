"""
PPT Generator Module

This module generates PowerPoint presentations from JSON data containing user information.
It creates slides with formatted sections for user details, interests, living characteristics,
appliance needs, and clothing insights.
"""

import json
import os
from typing import Dict, Tuple, List, Any, Optional

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT

# 颜色常量
COLORS = {
    'PRIMARY_GREEN': RGBColor(146, 208, 80),   # 主要绿色边框
    'SECONDARY_GREEN': RGBColor(178, 210, 52), # 次要绿色边框
    'BLACK': RGBColor(0, 0, 0),                # 黑色文字
    'GREEN_TEXT': RGBColor(101, 175, 69),      # 绿色文字
    'GRAY': RGBColor(200, 200, 200),           # 灰色文字
    'LIGHT_GRAY': RGBColor(240, 240, 240),     # 浅灰色背景
    'BORDER_GRAY': RGBColor(220, 220, 220),    # 灰色边框
    'WHITE': RGBColor(255, 255, 255),          # 白色文字
}

# 字体和布局常量
FONT_NAME = '微软雅黑'
BULLET_SYMBOL = "●"

def extract_json_info(user_data: Dict[str, Any]) -> Tuple[Dict[str, str], Dict[str, str], Dict[str, str], Dict[str, str], Dict[str, str]]:
    """
    从JSON数据中提取用户信息，并组织成PPT各区域所需的格式
    
    Args:
        user_data: 包含用户信息的JSON数据字典
        
    Returns:
        
    """
    page = {}

    page1_tags = {}
    page2_tags = {}
    page3_tags = {}

    # Page1区域数据字典
    basic_info = {}
    interests = {}
    living_features = {}
    appliance_needs = {}
    clothing_insights = {}

    #Page2区域数据字典
    #品类选择考虑
    page2_left = {}
    #核心关注因素展开
    page2_center = {}
    #综合
    page2_right = {}
    
    
    #Page3区域数据字典
    #产品使用习惯
    page3_left = {}
    page3_left_2 = {}
    #产品使用评价&痛点&爽点
    page3_right_top = {}
    #Leader品牌&产品升级机会
    page3_right_bom = {}

    # 获取常用的数据节点，避免重复查询
    section_A = user_data.get("A. 年轻人生活状态探讨", {})
    section_B = user_data.get("B. 波轮洗衣机购买全链路还原", {})
    section_C = user_data.get("C. 年轻人衣物洗干全场景洞察", {})
    section_D = user_data.get("D. 产品使用及满意度评价", {})
    section_F = user_data.get("F. 探索 Leader 品牌 & 产品升级机会", {})

    


    
    basic_info_section = section_A.get("3、基本信息", {})
    product_info_section = section_A.get("2、产品信息", {})
    interests_section = section_A.get("4、兴趣爱好", {})
    living_section = section_A.get("8、居住特征", {})




    page1_tags['tag'] = "用户标签：年轻单身女性、职场新人、月光族、租房独居青年"
    page2_tags['tag'] = "购买标签：低预算、选知名品牌、选购时长半个月、多平台对比"
    page3_tags['tag'] = "使用标签：关注消毒、程序自定义、观察洗衣过程"
    
    page2_right['信息搜索路径'] = section_B.get("14、信息搜索路径还原", {}).get("信息搜索路径", "")
    page2_right['品牌对比'] = section_B.get("12、购买过程", {}).get("品牌对比", "")
    page2_right['最终产品购买原因'] = section_B.get("12、购买过程", {}).get("最终购买原因", "")
    page2_right['排除其他品牌原因'] = section_B.get("12、购买过程", {}).get("品牌态度", "")
    page2_right['购买障碍'] = section_B.get("13、购买决策", {}).get("购买障碍", "")
    page2_right['购买渠道'] = section_B.get("13、购买决策", {}).get("购买渠道", "")
    page2_right['决策时长'] = section_B.get("14、信息搜索路径还原", {}).get("决策时长", "")
    page2_right['传统电商类内容关注'] = section_B.get("14、信息搜索路径还原", {}).get("传统电商类内容关注", "")
    page2_right['电商详情页偏好'] = section_B.get("14、信息搜索路径还原", {}).get("", "电商详情页偏好")
    page2_right['社媒内容关注'] = section_B.get("14、信息搜索路径还原", {}).get("社媒内容关注", "")
    page2_right['KOL关注'] = section_B.get("14、信息搜索路径还原", {}).get("KOL关注", "")
    page2_right['其他信息内容关注'] = ""

    page2_left['购买动机'] = section_A.get("11、购买产品的背景及动机", {}).get("购买动机", "")
    page2_left['过往经验'] = section_A.get("11、购买产品的背景及动机", {}).get("过往使用经验及品类态度", "")
    page2_left['品类态度'] = ""
    page2_left['波轮定位'] = section_A.get("11、购买产品的背景及动机", {}).get("波轮洗衣机的定位", "")
    page2_left['预算'] = section_B.get("12、购买过程", {}).get("价格", "")
    
    page2_center['核心关注因素排序'] = section_B.get("13、购买决策", {}).get("购买决策因素（KBF）", "")
    page2_center['外观'] = section_B.get("12、购买过程", {}).get("外观", "")
    page2_center['屏幕与操作方式'] = section_B.get("12、购买过程", {}).get("屏幕与操控方式", "")
    page2_center['尺寸'] = section_B.get("12、购买过程", {}).get("尺寸", "")
    page2_center['容量'] = section_B.get("12、购买过程", {}).get("容量", "")
    page2_center['洗涤程序'] = section_B.get("12、购买过程", {}).get("", "洗涤程序")
    page2_center['功能键'] = section_B.get("12、购买过程", {}).get("", "功能键")
    page2_center['防缠绕'] = section_B.get("12、购买过程", {}).get("", "防缠绕")
    page2_center['洗净比'] = section_B.get("12、购买过程", {}).get("", "洗净比")
    page2_center['水流'] = section_B.get("12、购买过程", {}).get("水流", "")
    page2_center['噪音'] = section_B.get("12、购买过程", {}).get("噪音", "")
    page2_center['耗电量、耗水量'] = section_B.get("12、购买过程", {}).get("", "能耗等级、耗电量、耗水量")
    page2_center['除菌、除螨'] = section_B.get("12、购买过程", {}).get("除菌/除螨方式", "")
    page2_center['电机'] = section_B.get("12、购买过程", {}).get("电机", "")
    page2_center['转速'] = section_B.get("12、购买过程", {}).get("转速", "")
    page2_center['筒自洁'] = section_B.get("12、购买过程", {}).get("筒自洁", "")
    page2_center['智能化'] = section_B.get("12、购买过程", {}).get("智能化", "")
    page2_center['促销'] = section_B.get("12、购买过程", {}).get("促销活动", "")
    page2_center['售后/质保'] = section_B.get("12、购买过程", {}).get("售后/质保", "")


    page3_left['洗衣频率'] = section_C.get("16、衣物洗护习惯", {}).get("洗衣频率", "")
    page3_left['洗衣时间'] = section_C.get("16、衣物洗护习惯", {}).get("洗衣时间", "")
    page3_left['攒脏衣地点'] = section_C.get("17、机洗习惯", {}).get("存放脏衣服", "")
    page3_left['洗衣分类习惯'] = section_C.get("17、机洗习惯", {}).get("分-衣物分类习惯及痛点", "")
    page3_left['手洗习惯'] = section_C.get("16、衣物洗护习惯", {}).get("洗护方式-手洗/干洗", "")
    page3_left['外送习惯'] = section_C.get("16、衣物洗护习惯", {}).get("送外洗衣习惯", "")
    page3_left['预洗习惯'] = section_C.get("17、机洗习惯", {}).get("预洗-预洗习惯及痛点", "")
    page3_left['浸泡习惯'] = section_C.get("17、机洗习惯", {}).get("预洗-预洗习惯及痛点", "")
    page3_left['洗涤剂使用习惯（洗涤剂+消毒剂+柔顺剂）'] = section_C.get("17、机洗习惯", {}).get("投-洗涤剂投放习惯及痛点", "")

    page3_left_2['最常用程序'] = section_C.get("17、机洗习惯", {}).get("选-程序选择习惯及投放", "")
    page3_left_2['手动调参数习惯'] = section_C.get("17、机洗习惯", {}).get("手动调参数习惯", "")
    page3_left_2['智能操控习惯'] = section_C.get("17、机洗习惯", {}).get("智能操控习惯", "")

    page3_left['程序选择习惯'] = page3_left_2

    page3_left['单脱水习惯'] = section_C.get("17、机洗习惯", {}).get("脱-单独脱水习惯及痛点", "")
    page3_left['晾晒习惯'] = section_C.get("17、机洗习惯", {}).get("干-干衣方式", "")
    page3_left['机器清洁习惯'] = section_C.get("17、机洗习惯", {}).get("清-洗衣机/干衣机清洗及维护习惯及痛点", "")
    page3_left['典型使用场景'] = "待补充..."

    page3_right_top['综合评分'] = section_D.get("18、产品使用评价", {}).get("综合打分", "")
    page3_right_top['痛点'] = section_D.get("19、深度评价", {}).get("产品使用痛点", "")
    page3_right_top['爽点'] = section_D.get("19、深度评价", {}).get("产品使用爽点", "")
    page3_right_top['推荐度'] = section_D.get("20、产品推荐度", {}).get("推荐度", "")+"&"+section_D.get("20、产品推荐度", {}).get("如何推荐", "")

    page3_right_bom['Leader品牌印象'] = section_F.get("23、Leader品牌&产品升级机会", {}).get("Leader用户-品牌印象", "")
    page3_right_bom['年轻品牌特征'] = section_F.get("23、Leader品牌&产品升级机会", {}).get("所有用户-年轻品牌特征", "")
    page3_right_bom['波轮品牌期待'] = section_F.get("23、Leader品牌&产品升级机会", {}).get("所有用户-品牌期待", "")

    page3_right_bom['波轮产品期待'] = section_F.get("23、Leader品牌&产品升级机会", {}).get("所有用户-产品期待", "")
    page3_right_bom['波轮价格期待'] = section_F.get("23、Leader品牌&产品升级机会", {}).get("所有用户-价格期待", "")
    page3_right_bom['波轮营销期待'] = section_F.get("23、Leader品牌&产品升级机会", {}).get("所有用户-营销期待", "")



    # 用户基本信息
    age = basic_info_section.get('年龄', '')
    gender = basic_info_section.get('性别', '')
    basic_info['性别&年龄'] = f"{age}岁 {gender}" if age else gender
    
    basic_info['职业'] = basic_info_section.get("行业及职业", "")
    
    family_structure = basic_info_section.get("家庭结构", "")
    children_info = basic_info_section.get("孩子情况", "")
    basic_info['家庭结构'] = f"{family_structure}&{children_info}" if family_structure and children_info else family_structure or children_info
    
    basic_info['家庭年收入'] = basic_info_section.get("家庭年收入", "")
    
    brand = product_info_section.get("品牌", "")
    price = product_info_section.get("价格", "")
    basic_info['洗衣机品牌&价格'] = f"{brand}&{price}" if brand and price else brand or price
    
    basic_info['日常触媒习惯'] = section_A.get("10、触媒习惯", {}).get("触媒渠道", "")

    # 生活方式/兴趣爱好/消费观
    interests['工作节奏'] = interests_section.get('工作节奏', '')
    interests['生活重心'] = interests_section.get('生活重心', '')
    interests['兴趣爱好'] = interests_section.get('个人喜好', '')
    interests['消费观'] = section_A.get("9、家电需求", {}).get('家电消费观', '')

    # 住房情况
    living_features['房子情况'] = living_section.get('房子情况', '')
    living_features['装修花费'] = ''  # 保持原有的空值
    living_features['家装风格'] = living_section.get('家装风格', '')

    # 家电盘点
    appliance_needs['其他洗涤设备'] = living_section.get('洗涤设备', '')
    appliance_needs['大家电盘点'] = living_section.get('家电盘点', '')
    appliance_needs['家电智能化程度'] = living_section.get('家电智能化程度', '')

    # 衣物洞察
    daily_habits_section = section_C.get("15、日常穿着习惯（家访前小问卷数据导入）", {})
    clothing_insights['日常衣物类型'] = daily_habits_section.get('衣物类型', '')
    clothing_insights['兴趣爱好相关特殊衣物洞察'] = section_A.get("5、特殊衣物", {}).get('特殊衣物洞察', '')
    clothing_insights['衣物价格范围'] = daily_habits_section.get('服饰购买价格上下限', '')

    page['page1']=[page1_tags,basic_info, interests, living_features, appliance_needs, clothing_insights]
    page['page2']=[page2_tags,page2_left,page2_center,page2_right]
    page['page3']=[page3_tags,page3_left,page3_right_top,page3_right_bom]

    return page


def create_slide(prs: Presentation, pagedata: Dict[str, Any], user_id: str) -> None:
    """
    创建单个用户的PPT幻灯片
    
    Args:
        prs: PowerPoint演示文稿对象
        user_data: 用户数据字典
        user_id: 用户ID，用于显示在幻灯片标题
    """
    # 添加一个空白幻灯片
    slide_layout = prs.slide_layouts[6]  # 使用空白布局
    slide = prs.slides.add_slide(slide_layout)
    

    page1_tags, basic_info, interests, living_features, appliance_needs, clothing_insights= pagedata
    
    # 添加标题
    add_title(slide, user_id)
    

    # 添加用户标签行
    add_subtitle(slide,page1_tags['tag'])
    
    # 创建五个主要区域 - 调整位置和大小以适应单页
    section_configs = [
        ("用户基本信息", 0.3, 1.2, 3.5, 6.1, basic_info),
        ("兴趣爱好", 3.9, 1.2, 3.3, 3.4, interests),
        ("居住特征", 7.3, 1.2, 2.8, 3.4, living_features),
        ("家电需求", 10.2, 1.2, 2.9, 3.4, appliance_needs),
        ("衣物洞察", 3.9, 4.7, 9.2, 2.6, clothing_insights)
    ]
    
    for title, left, top, width, height, data in section_configs:
        create_section(slide, title, left, top, width, height, data)


def create_slide2(prs: Presentation, pagedata: Dict[str, Any], location_id: str) -> None:
    """
    创建洗衣机购买考虑因素的PPT幻灯片
    
    Args:
        prs: PowerPoint演示文稿对象
        data: 数据字典
        location_id: 位置ID，用于显示在幻灯片标题
    """
    # 添加一个空白幻灯片
    slide_layout = prs.slide_layouts[6]  # 使用空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 提取洗衣机购买相关信息
    page2_tags,page2_left,page2_center,page2_right=pagedata
    
    # 添加标题
    add_title(slide, location_id)
    
    # 添加副标题
    add_subtitle(slide,page2_tags['tag'])
    
    # 创建三个主要区域
    # 左侧：洗衣机购买动机
    create_page2_left_section(slide, 0.3, 1.2, 2.8, 5.2, page2_left)
    
    # 中间：洗衣机选购关注因素
    create_page2_center_section(slide, 3.2, 1.7, 3.7, 5.7, page2_center)
    
    # 右侧：洗衣机选购对比
    create_page2_right_section(slide, 6.8, 1.7, 6.2, 5.7, page2_right)

def create_slide3(prs: Presentation, pagedata: Dict[str, Any], location_id: str) -> None:
    """
    创建洗衣机购买考虑因素的PPT幻灯片
    
    Args:
        prs: PowerPoint演示文稿对象
        data: 数据字典
        location_id: 位置ID，用于显示在幻灯片标题
    """
    # 添加一个空白幻灯片
    slide_layout = prs.slide_layouts[6]  # 使用空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 提取洗衣机购买相关信息
    page3_tags,page3_left,page3_right_top,page3_right_bom=pagedata
    
    # 添加标题
    add_title(slide, location_id)
    
    # 添加副标题
    add_subtitle(slide,page3_tags['tag'])
    
    # 创建三个主要区域
    # 左侧
    create_page3_left_section(slide, 0.3, 1.2, 6.8, 5.0, page3_left)
    
    # 右上
    create_page3_right_top_section(slide, 7.6, 1.2, 5.2, 5.0, page3_right_top)
    
    # 右下
    create_page3_right_bom_section(slide, 7.6, 3.9, 5.2, 5.0, page3_right_bom)


def create_slide4(prs: Presentation, location_id: str) -> None:
    """
    创建洗衣机购买考虑因素的PPT幻灯片
    
    Args:
        prs: PowerPoint演示文稿对象
        data: 数据字典
        location_id: 位置ID，用于显示在幻灯片标题
    """
    # 添加一个空白幻灯片
    slide_layout = prs.slide_layouts[6]  # 使用空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    
    # 添加标题
    title = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(12), Inches(0.5))
    title_frame = title.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = f"【{location_id}】"
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['GREEN_TEXT']
    title_para.font.name = FONT_NAME

    # 添加值
    run = title_para.add_run()
    run.text = f"用户照片"
    run.font.name = FONT_NAME
    run.font.size = Pt(24)
    run.font.color.rgb = COLORS['BLACK']
    run.font.bold = True
    
    # 创建区域
    section_configs = [
        ("小区环境", 0.5, 0.7, 6, 1.7),
        ("兴趣爱好", 6.7, 0.7, 6, 1.7),
        ("家居环境", 0.5, 2.6, 12.2, 1.9),
        ("产品相关", 0.5, 4.7, 6, 2.5),
        ("入户发现", 6.7, 4.7, 6, 2.5)
    ]
    
    for title, left, top, width, height in section_configs:
        create_page4_section(slide, title, left, top, width, height)
    
def create_page4_section(slide, title: str, left: float, top: float, width: float, height: float) -> None:
    # 添加上左面区域边框
    border = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(2.0), Inches(0.3))
    # 自定义颜色填充,178,210,52
    border.fill.solid()
    border.fill.fore_color.rgb = RGBColor(139, 197, 63)
    # 在border1里添加文字，非单独的文本框
    text_frame = border.text_frame  # 获取文本框对象
    # ===== Step 2: 清空默认段落（可选） =====
    text_frame.clear()
    # ===== Step 3: 添加段落并设置文字内容 =====
    paragraph = text_frame.paragraphs[0]
    paragraph.text = title  # 设置你的文本
    # ===== Step 4: 调整文字格式（按需定制） =====
    paragraph.alignment = PP_ALIGN.CENTER  # 居中对齐
    run = paragraph.runs[0]  # 获取文本段落的Run对象
    run.font.size = Pt(18)   # 字体大小
    run.font.name = FONT_NAME
    # 字体颜色白色
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.bold = True     # 加粗


    # 添加区域边框
    border = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    border.fill.background()  # 透明填充
    border.line.color.rgb = COLORS['PRIMARY_GREEN']
    border.line.width = Pt(1.5)

def create_page3_left_section(slide, left: float, top: float, width: float, height: float, data: Dict[str, str]) -> None:
    """
    创建洗衣机购买动机区域
    
    Args:
        slide: 幻灯片对象
        left: 左边距（英寸）
        top: 上边距（英寸）
        width: 宽度（英寸）
        height: 高度（英寸）
        data: 区域内容数据字典
    """

    # 添加上左面区域边框
    border1 = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height-4.6))
    # 自定义颜色填充,178,210,52
    border1.fill.solid()
    border1.fill.fore_color.rgb = RGBColor(139, 197, 63)
    # 在border1里添加文字，非单独的文本框
    text_frame = border1.text_frame  # 获取文本框对象
    # ===== Step 2: 清空默认段落（可选） =====
    text_frame.clear()
    # ===== Step 3: 添加段落并设置文字内容 =====
    paragraph = text_frame.paragraphs[0]
    paragraph.text = "产品使用习惯"  # 设置你的文本
    # ===== Step 4: 调整文字格式（按需定制） =====
    paragraph.alignment = PP_ALIGN.CENTER  # 居中对齐
    run = paragraph.runs[0]  # 获取文本段落的Run对象
    run.font.size = Pt(16)   # 字体大小
    run.font.name = FONT_NAME
    # 字体颜色白色
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.bold = True     # 加粗


    # 添加区域边框
    border = slide.shapes.add_shape(1, Inches(left), Inches(top+0.4), Inches(width), Inches(height+0.5))
    border.fill.background()  # 透明填充
    border.line.color.rgb = COLORS['PRIMARY_GREEN']
    border.line.width = Pt(1.5)
    
    # 添加内容文本框
    content_box = slide.shapes.add_textbox(
        Inches(left + 0.1),
        Inches(top + 0.5),
        Inches(width - 0.2),
        Inches(height+0.3)
    )

    
    
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for key, value in data.items():
        if key == list(data.keys())[0]:
            para = content_frame.paragraphs[0]
        else:
            para = content_frame.add_paragraph()            
        

        para.text = f"{BULLET_SYMBOL} {key}："
        para.font.size = Pt(12)
        para.font.bold = True
        para.font.name = FONT_NAME
        para.font.color.rgb = COLORS['BLACK']
        para.space_after = Pt(8)
        
        

        if key == "程序选择习惯":
            for key2,val2 in value.items():
                para2 = content_frame.add_paragraph()
                # 设置段落格式
                para2.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                para2.space_after = Pt(1)
                para2.line_spacing = 1.0
                
                para2.text = f"{key2}： "
                para2.level = 1

                
                # 设置字体
                para2.font.name = FONT_NAME
                para2.font.size = Pt(12)
                para2.font.color.rgb = COLORS['GREEN_TEXT']
                para2.font.bold = True
                
                # 添加值
                run = para2.add_run()
                run.text = f"{val2}"
                run.font.name = FONT_NAME
                run.font.size = Pt(12)
                run.font.color.rgb = COLORS['BLACK']
                run.font.bold = False
        else:
            
            # 添加值
            run = para.add_run()
            run.text = f"{value}"
            run.font.name = FONT_NAME
            run.font.size = Pt(12)
            run.font.color.rgb = COLORS['BLACK']
            run.font.bold = False

            if key == "典型使用场景":
                para.font.color.rgb = RGBColor(31, 73, 125)
                run.font.color.rgb = RGBColor(250, 0, 0)


def create_page3_right_top_section(slide, left: float, top: float, width: float, height: float, data: Dict[str, str]) -> None:
    """
    创建洗衣机购买动机区域
    
    Args:
        slide: 幻灯片对象
        left: 左边距（英寸）
        top: 上边距（英寸）
        width: 宽度（英寸）
        height: 高度（英寸）
        data: 区域内容数据字典
    """

    # 添加上左面区域边框
    border1 = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height-4.6))
    # 自定义颜色填充,178,210,52
    border1.fill.solid()
    border1.fill.fore_color.rgb = RGBColor(139, 197, 63)
    # 在border1里添加文字，非单独的文本框
    text_frame = border1.text_frame  # 获取文本框对象
    # ===== Step 2: 清空默认段落（可选） =====
    text_frame.clear()
    # ===== Step 3: 添加段落并设置文字内容 =====
    paragraph = text_frame.paragraphs[0]
    paragraph.text = "产品使用评价&痛点&爽点"  # 设置你的文本
    # ===== Step 4: 调整文字格式（按需定制） =====
    paragraph.alignment = PP_ALIGN.CENTER  # 居中对齐
    run = paragraph.runs[0]  # 获取文本段落的Run对象
    run.font.size = Pt(16)   # 字体大小
    run.font.name = FONT_NAME
    # 字体颜色白色
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.bold = True     # 加粗


    # 添加区域边框
    border = slide.shapes.add_shape(1, Inches(left), Inches(top+0.4), Inches(width), Inches(height-2.8))
    border.fill.background()  # 透明填充
    border.line.color.rgb = COLORS['PRIMARY_GREEN']
    border.line.width = Pt(1.5)
    
    # 添加内容文本框
    content_box = slide.shapes.add_textbox(
        Inches(left + 0.1),
        Inches(top + 0.5),
        Inches(width - 0.2),
        Inches(height-2.8)
    )

    
    
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for key, value in data.items():
        if key == list(data.keys())[0]:
            para = content_frame.paragraphs[0]
        else:
            para = content_frame.add_paragraph()            
        

        para.text = f"{BULLET_SYMBOL} {key}："
        para.font.size = Pt(12)
        para.font.bold = True
        para.font.name = FONT_NAME
        para.font.color.rgb = COLORS['BLACK']
        para.space_after = Pt(6)
        
            
        # 添加值
        run = para.add_run()
        run.text = f"{value}"
        run.font.name = FONT_NAME
        run.font.size = Pt(12)
        run.font.color.rgb = COLORS['BLACK']
        run.font.bold = False

          
        
def create_page3_right_bom_section(slide, left: float, top: float, width: float, height: float, data: Dict[str, str]) -> None:
    """
    创建洗衣机购买动机区域
    
    Args:
        slide: 幻灯片对象
        left: 左边距（英寸）
        top: 上边距（英寸）
        width: 宽度（英寸）
        height: 高度（英寸）
        data: 区域内容数据字典
    """

    # 添加上左面区域边框
    border1 = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height-4.6))
    # 自定义颜色填充,178,210,52
    border1.fill.solid()
    border1.fill.fore_color.rgb = RGBColor(139, 197, 63)
    # 在border1里添加文字，非单独的文本框
    text_frame = border1.text_frame  # 获取文本框对象
    # ===== Step 2: 清空默认段落（可选） =====
    text_frame.clear()
    # ===== Step 3: 添加段落并设置文字内容 =====
    paragraph = text_frame.paragraphs[0]
    paragraph.text = "产品使用评价&痛点&爽点"  # 设置你的文本
    # ===== Step 4: 调整文字格式（按需定制） =====
    paragraph.alignment = PP_ALIGN.CENTER  # 居中对齐
    run = paragraph.runs[0]  # 获取文本段落的Run对象
    run.font.size = Pt(16)   # 字体大小
    run.font.name = FONT_NAME
    # 字体颜色白色
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.bold = True     # 加粗


    # 添加区域边框
    border = slide.shapes.add_shape(1, Inches(left), Inches(top+0.4), Inches(width), Inches(height-2.2))
    border.fill.background()  # 透明填充
    border.line.color.rgb = COLORS['PRIMARY_GREEN']
    border.line.width = Pt(1.5)
    
    # 添加内容文本框
    content_box = slide.shapes.add_textbox(
        Inches(left + 0.1),
        Inches(top + 0.5),
        Inches(width - 0.2),
        Inches(height-2.4)
    )

    
    
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for key, value in data.items():
        if key == list(data.keys())[0]:
            para = content_frame.paragraphs[0]
        else:
            para = content_frame.add_paragraph()            
        

        para.text = f"{BULLET_SYMBOL} {key}："
        para.font.size = Pt(12)
        para.font.bold = True
        para.font.name = FONT_NAME
        para.font.color.rgb = COLORS['BLACK']
        para.space_after = Pt(8)
        
            
        # 添加值
        run = para.add_run()
        run.text = f"{value}"
        run.font.name = FONT_NAME
        run.font.size = Pt(12)
        run.font.color.rgb = COLORS['BLACK']
        run.font.bold = False





def create_page2_left_section(slide, left: float, top: float, width: float, height: float, data: Dict[str, str]) -> None:
    """
    创建洗衣机购买动机区域
    
    Args:
        slide: 幻灯片对象
        left: 左边距（英寸）
        top: 上边距（英寸）
        width: 宽度（英寸）
        height: 高度（英寸）
        data: 区域内容数据字典
    """

    # 添加上左面区域边框
    border1 = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height-4.8))
    # 自定义颜色填充,178,210,52
    border1.fill.solid()
    border1.fill.fore_color.rgb = RGBColor(178, 210, 52)
    # 在border1里添加文字，非单独的文本框
    text_frame = border1.text_frame  # 获取文本框对象
    # ===== Step 2: 清空默认段落（可选） =====
    text_frame.clear()  # 
    # ===== Step 3: 添加段落并设置文字内容 =====
    paragraph = text_frame.paragraphs[0]
    paragraph.text = "洗衣机购买动机"  # 设置你的文本
    # ===== Step 4: 调整文字格式（按需定制） =====
    paragraph.alignment = PP_ALIGN.CENTER  # 居中对齐
    run = paragraph.runs[0]  # 获取文本段落的Run对象
    run.font.size = Pt(16)   # 字体大小
    run.font.name = FONT_NAME
    # 字体颜色白色
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.bold = True     # 加粗


    # 添加上右面区域边框
    border2 = slide.shapes.add_shape(1, Inches(left+width+0.1), Inches(top), Inches(width*3.5), Inches(height-4.8))
    # 自定义颜色填充,178,210,52
    border2.fill.solid()
    border2.fill.fore_color.rgb = RGBColor(139, 197, 63)
    # 在border1里添加文字，非单独的文本框
    text_frame = border2.text_frame  # 获取文本框对象
    # ===== Step 2: 清空默认段落（可选） =====
    text_frame.clear()  # 
    # ===== Step 3: 添加段落并设置文字内容 =====
    paragraph = text_frame.paragraphs[0]
    paragraph.text = "洗衣机选购对比"  # 设置你的文本
    # ===== Step 4: 调整文字格式（按需定制） =====
    paragraph.alignment = PP_ALIGN.CENTER  # 居中对齐
    run = paragraph.runs[0]  # 获取文本段落的Run对象
    run.font.size = Pt(16)   # 字体大小
    run.font.name = FONT_NAME
    # 字体颜色白色
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.bold = True     # 加粗


    # 添加区域边框
    border = slide.shapes.add_shape(1, Inches(left), Inches(top+0.5), Inches(width), Inches(height+0.5))
    border.fill.background()  # 透明填充
    border.line.color.rgb = COLORS['PRIMARY_GREEN']
    border.line.width = Pt(1.5)
    
    # 添加内容文本框
    content_box = slide.shapes.add_textbox(
        Inches(left + 0.1),
        Inches(top + 0.6),
        Inches(width - 0.2),
        Inches(height+0.3)
    )

    
    
    content_frame = content_box.text_frame
    content_frame.word_wrap = True


    # 添加购买动机标题
    para = content_frame.paragraphs[0]

    para.text = f"{BULLET_SYMBOL} 购买动机"
    para.font.size = Pt(16)
    para.font.bold = True
    para.font.name = FONT_NAME
    para.font.color.rgb = COLORS['BLACK']
    para.space_after = Pt(10)
    
    # 添加购买动机内容
    para = content_frame.add_paragraph()
    para.text = data['购买动机']
    para.font.size = Pt(12)
    para.font.name = FONT_NAME
    para.level = 0
    para.space_after = Pt(40)
    
    # 添加品类选择考虑标题
    para = content_frame.add_paragraph()
    para.text = f"{BULLET_SYMBOL} 品类选择考虑"
    para.font.size = Pt(16)
    para.font.bold = True
    para.font.name = FONT_NAME
    para.font.color.rgb = COLORS['BLACK']
    para.space_after = Pt(10)
    
    data.pop('购买动机', None)

    # 循环字典data


    for key, value in data.items():
        para = content_frame.add_paragraph()
        # 设置段落格式
        para.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        para.space_after = Pt(1)
        para.line_spacing = 1.0
        
        para.text = f"{key}： "

        # 设置段落项目符号
        para.level = 0
        para.level_text = BULLET_SYMBOL
        
        # 设置字体
        para.font.name = FONT_NAME
        para.font.size = Pt(12)
        para.font.color.rgb = COLORS['GREEN_TEXT']
        para.font.bold = True
        
        # 添加值
        run = para.add_run()
        run.text = f"{value}"
        run.font.name = FONT_NAME
        run.font.size = Pt(12)
        run.font.color.rgb = COLORS['BLACK']
        run.font.bold = False


def create_page2_center_section(slide, left: float, top: float, width: float, height: float, data: Dict[str, Any]) -> None:
    """
    创建洗衣机选购关注因素区域
    
    Args:
        slide: 幻灯片对象
        left: 左边距（英寸）
        top: 上边距（英寸）
        width: 宽度（英寸）
        height: 高度（英寸）
        data: 区域内容数据字典
    """
    # 添加区域边框
    border = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    border.fill.background()  # 透明填充
    border.line.color.rgb = COLORS['PRIMARY_GREEN']
    border.line.width = Pt(1.5)
    
    # 添加内容文本框
    content_box = slide.shapes.add_textbox(
        Inches(left + 0.1),
        Inches(top+0.1),
        Inches(width - 0.2),
        Inches(height - 0.4)
    )

    
    
    content_frame = content_box.text_frame
    content_frame.word_wrap = True


    # 添加核心关注因素排序标题
    para = content_frame.paragraphs[0]

    para.text = f"{BULLET_SYMBOL} 核心关注因素排序"
    para.font.size = Pt(16)
    para.font.bold = True
    para.font.name = FONT_NAME
    para.font.color.rgb = COLORS['BLACK']
    para.space_after = Pt(10)
    
    # 添加核心关注因素排序内容
    para = content_frame.add_paragraph()
    para.text = data['核心关注因素排序']
    para.font.size = Pt(12)
    para.font.name = FONT_NAME
    para.level = 0
    para.space_after = Pt(20)
    
    # 添加品类选择考虑标题
    para = content_frame.add_paragraph()
    para.text = f"{BULLET_SYMBOL} 核心关注因素展开"
    para.font.size = Pt(16)
    para.font.bold = True
    para.font.name = FONT_NAME
    para.font.color.rgb = COLORS['BLACK']
    para.space_after = Pt(10)
    
    data.pop('购买动机', None)

    # 循环字典data


    for key, value in data.items():
        para = content_frame.add_paragraph()
        # 设置段落格式
        para.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        para.space_after = Pt(1)
        para.line_spacing = 1.0
        
        para.text = f"{key}： "

        # 设置段落项目符号
        para.level = 0
        para.level_text = BULLET_SYMBOL
        
        # 设置字体
        para.font.name = FONT_NAME
        para.font.size = Pt(10)
        para.font.color.rgb = COLORS['GREEN_TEXT']
        para.font.bold = True
        
        # 添加值
        run = para.add_run()
        run.text = f"{value}"
        run.font.name = FONT_NAME
        run.font.size = Pt(10)
        run.font.color.rgb = COLORS['BLACK']
        run.font.bold = False


def create_page2_right_section(slide, left: float, top: float, width: float, height: float, data: Dict[str, Any]) -> None:
    """
    创建洗衣机选购对比区域
    
    Args:
        slide: 幻灯片对象
        left: 左边距（英寸）
        top: 上边距（英寸）
        width: 宽度（英寸）
        height: 高度（英寸）
        data: 区域内容数据字典
    """
    
    border = slide.shapes.add_shape(1, Inches(left+0.2), Inches(top), Inches(width-0.2), Inches(height))
    border.fill.background()  # 透明填充
    border.line.color.rgb = COLORS['PRIMARY_GREEN']
    border.line.width = Pt(1.5)
    
    # 添加内容文本框
    content_box = slide.shapes.add_textbox(
        Inches(left + 0.4),
        Inches(top+0.1),
        Inches(width - 0.5),
        Inches(height - 0.4)
    )

    
    
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for key, value in data.items():
        if key == list(data.keys())[0]:
            para = content_frame.paragraphs[0]
        else:
            para = content_frame.add_paragraph()            

        para.text = f"{BULLET_SYMBOL} {key}："
        para.font.size = Pt(14)
        para.font.bold = True
        para.font.name = FONT_NAME
        para.font.color.rgb = COLORS['BLACK']
        para.space_after = Pt(10)
        # 添加值
        run = para.add_run()
        run.text = value
        run.font.name = FONT_NAME
        run.font.size = Pt(12)
        run.font.color.rgb = COLORS['BLACK']
        run.font.bold = False


def add_title(slide, user_id: str) -> None:
    """
    添加幻灯片标题
    
    Args:
        slide: 幻灯片对象
        user_id: 用户ID
    """
    title = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(12), Inches(0.5))
    title_frame = title.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = f"【{user_id}】"
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['GREEN_TEXT']
    title_para.font.name = FONT_NAME


def add_subtitle(slide,tags) -> None:
    """
    添加幻灯片副标题（用户标签）
    
    Args:
        slide: 幻灯片对象
    """
    subtitle = slide.shapes.add_textbox(Inches(0.2), Inches(0.6), Inches(10), Inches(0.5))
    subtitle_frame = subtitle.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = tags
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.bold = True
    subtitle_para.font.name = FONT_NAME

def create_section(slide, title: str, left: float, top: float, width: float, height: float, data: Dict[str, str]) -> None:
    """
    创建幻灯片中的单个区域
    
    Args:
        slide: 幻灯片对象
        title: 区域标题
        left: 左边距（英寸）
        top: 上边距（英寸）
        width: 宽度（英寸）
        height: 高度（英寸）
        data: 区域内容数据字典
    """
    # 添加区域边框
    border = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    border.fill.background()  # 透明填充
    
    # 设置边框颜色
    border.line.color.rgb = COLORS['PRIMARY_GREEN'] if title == "用户基本信息" else COLORS['SECONDARY_GREEN']
    border.line.width = Pt(1.5)
    
    # 添加区域标题
    add_section_title(slide, title, left, top, width)
    
    # 设置内容区域的垂直位置
    top_var = 3.2 if title == "用户基本信息" else 0.2
    
    # 添加内容文本框
    content_box = add_content_box(slide, left, top+top_var, width, height )
    content_frame = content_box.text_frame
    
    # 为用户基本信息添加照片区域
    if title == "用户基本信息":
        add_photo_placeholder(slide, left, top)
    
    # 添加数据内容
    add_data_content(content_frame, data)


def add_section_title(slide, title: str, left: float, top: float, width: float) -> None:
    """
    添加区域标题
    
    Args:
        slide: 幻灯片对象
        title: 标题文本
        left: 左边距（英寸）
        top: 上边距（英寸）
        width: 宽度（英寸）
    """
    title_box = slide.shapes.add_textbox(
        Inches(left + 0.2), 
        Inches(top + 0.1), 
        Inches(width - 0.4), 
        Inches(0.4)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = f"{BULLET_SYMBOL} {title}"
    title_para.font.size = Pt(16)
    title_para.font.bold = True
    title_para.font.name = FONT_NAME


def add_content_box(slide, left: float, top: float, width: float, height: float):
    """
    添加内容文本框
    
    Args:
        slide: 幻灯片对象
        left: 左边距（英寸）
        top: 上边距（英寸）
        width: 宽度（英寸）
        height: 高度（英寸）
        
    Returns:
        创建的文本框对象
    """
    content_box = slide.shapes.add_textbox(
        Inches(left + 0.25),
        Inches(top + 0.05),
        Inches(width - 0.5),
        Inches(height - 1.3)
    )
    
    content_frame = content_box.text_frame
    
    # 设置文本框属性
    content_frame.word_wrap = True
    content_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    content_frame.vertical_anchor = MSO_ANCHOR.TOP
    
    return content_box


def add_photo_placeholder(slide, left: float, top: float) -> None:
    """
    添加照片占位区域
    
    Args:
        slide: 幻灯片对象
        left: 左边距（英寸）
        top: 上边距（英寸）
    """
    photo_bg = slide.shapes.add_shape(
        1,  # 矩形
        Inches(left + 0.3),
        Inches(top + 0.5),
        Inches(2.9),
        Inches(2.7)
    )
    photo_bg.fill.solid()
    photo_bg.fill.fore_color.rgb = COLORS['LIGHT_GRAY']
    photo_bg.line.color.rgb = COLORS['BORDER_GRAY']
    
    # 添加照片文本
    photo_text_frame = photo_bg.text_frame
    photo_text_frame.text = "用户本人照片"
    photo_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    photo_text_para = photo_text_frame.paragraphs[0]
    photo_text_para.font.size = Pt(18)
    photo_text_para.font.name = FONT_NAME
    photo_text_para.font.color.rgb = COLORS['WHITE']
    photo_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_data_content(content_frame, data: Dict[str, str]) -> None:
    """
    添加数据内容到文本框
    
    Args:
        content_frame: 文本框对象
        data: 数据字典，键为标题，值为内容
    """
    for key, value in data.items():
        # 创建新段落
        para = content_frame.add_paragraph()

        # 设置段落格式
        para.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        para.space_after = Pt(1)
        para.line_spacing = 1.0
        
        para.text = f"{key}： "

        # 设置段落项目符号
        para.level = 0
        para.level_text = BULLET_SYMBOL
        
        # 设置字体
        para.font.name = FONT_NAME
        para.font.size = Pt(12)
        para.font.color.rgb = COLORS['GREEN_TEXT']
        para.font.bold = True
        
        # 添加值
        run = para.add_run()
        run.text = f"{value}"
        run.font.name = FONT_NAME
        run.font.size = Pt(12)
        run.font.color.rgb = COLORS['BLACK']
        run.font.bold = False

def process_json_file(json_file_path: str, output_dir: str) -> str:
    """
    处理单个JSON文件并生成PPT
    
    Args:
        json_file_path: JSON文件路径
        output_dir: 输出目录路径
        
    Returns:
        生成的PPT文件路径
        
    Raises:
        FileNotFoundError: 当JSON文件不存在时
        json.JSONDecodeError: 当JSON文件格式不正确时
    """
    try:
        # 读取JSON文件
        with open(json_file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # 创建新的PPT
        prs = Presentation()
        
        # 设置幻灯片大小为16:9
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        data = dict(sorted(data.items(), key=lambda x: (x[0].split('\\')[0], x[1]['A. 年轻人生活状态探讨']['1、受访者']['序号'])))
        

        # 处理每个用户数据
        for user_id, user_data in data.items():
            # 从用户数据中提取标题信息
            ppt_title = extract_slide_title(user_data)
            # 提取用户信息
            page = extract_json_info(user_data)
            create_slide(prs, page['page1'], ppt_title)
            create_slide2(prs, page['page2'], ppt_title)
            create_slide3(prs, page['page3'], ppt_title)
            create_slide4(prs, ppt_title)
        
        # 获取文件名（不含扩展名）
        file_name = os.path.splitext(os.path.basename(json_file_path))[0]
        
        # 保存PPT - 使用原始文件名来确保唯一性
        output_path = os.path.join(output_dir, f"{file_name}_new.pptx")
        prs.save(output_path)
        print(f"已生成PPT: {output_path}")
        
        return output_path
    
    except FileNotFoundError:
        print(f"错误: 找不到文件 {json_file_path}")
        raise
    except json.JSONDecodeError:
        print(f"错误: {json_file_path} 不是有效的JSON文件")
        raise
    except Exception as e:
        print(f"处理文件 {json_file_path} 时发生错误: {str(e)}")
        raise


def extract_slide_title(user_data: Dict[str, Any]) -> str:
    """
    从用户数据中提取幻灯片标题
    
    Args:
        user_data: 用户数据字典
        
    Returns:
        格式化的幻灯片标题
    """
    ppt_title_json = user_data.get("A. 年轻人生活状态探讨", {}).get("1、受访者", {})
    
    # 获取序号并补零
    sequence = ppt_title_json.get("序号", "")
    sequence = sequence.zfill(3) if sequence else ""
    
    city = ppt_title_json.get("城市", "")
    interviewee = ppt_title_json.get("受访者", "")
    
    # 组合标题
    return f"{sequence}{city}-{interviewee}"


def generate_ppt(input_dir: str = "data/input/json/ai", output_dir: str = "data/output") -> List[str]:
    """
    生成PPT的主函数
    
    Args:
        input_dir: 输入JSON文件目录路径，默认为"data/input/json/ai"
        output_dir: 输出PPT文件目录路径，默认为"data/output"
        
    Returns:
        生成的PPT文件路径列表
    """
    # 确保输出目录存在
    os.makedirs(output_dir, exist_ok=True)
    
    # 处理目录中的所有JSON文件
    generated_files = []
    
    try:
        # 获取目录中的所有文件
        files = os.listdir(input_dir)
        
        # 过滤出JSON文件
        json_files = [f for f in files if f.endswith('.json')]
        
        if not json_files:
            print(f"警告: 目录 {input_dir} 中没有找到JSON文件")
            return []
        
        # 处理每个JSON文件
        for file_name in json_files:
            try:
                json_file_path = os.path.join(input_dir, file_name)
                output_path = process_json_file(json_file_path, output_dir)
                generated_files.append(output_path)
            except Exception as e:
                print(f"处理文件 {file_name} 时发生错误: {str(e)}")
                # 继续处理其他文件
                continue
    
    except FileNotFoundError:
        print(f"错误: 目录 {input_dir} 不存在")
    except PermissionError:
        print(f"错误: 没有权限访问目录 {input_dir}")
    except Exception as e:
        print(f"生成PPT时发生错误: {str(e)}")
    
    return generated_files

if __name__ == "__main__":
    try:
        generated_files = generate_ppt(input_dir="data/input/json/test/ppt")
        print(f"共生成 {len(generated_files)} 个PPT文件")
    except Exception as e:
        print(f"程序执行过程中发生错误: {str(e)}")